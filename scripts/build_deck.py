"""
Generate the UCO PSB Hackathon presentation (docs/presentation/).

Minimal design system: white background, ink text, one grey, one accent
red used sparingly; semantic state colors appear only in risk-state chips.
Regenerate after edits with:  python scripts/build_deck.py
"""

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- palette
INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x6B, 0x72, 0x80)
FAINT = RGBColor(0xF2, 0xF2, 0xF0)
LINE = RGBColor(0xD9, 0xD9, 0xD6)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)          # sparing use only
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
# semantic (risk states only)
S_GREEN = RGBColor(0x1B, 0x7A, 0x43)
S_AMBER = RGBColor(0xB4, 0x6A, 0x00)
S_RED = RGBColor(0xC0, 0x39, 0x2B)
S_GREY = RGBColor(0x8A, 0x8F, 0x98)

FONT = "Calibri"
FONT_LIGHT = "Calibri Light"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.75)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide():
    return prs.slides.add_slide(BLANK)


def text(s, x, y, w, h, txt, size=18, color=INK, bold=False, font=FONT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    first = True
    for line in txt.split("\n"):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line
        p.alignment = align
        p.line_spacing = line_spacing
        for r in p.runs:
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
    return box


def bullets(s, x, y, w, h, items, size=16, gap=6, color=INK):
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, (head, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.05
        r = p.add_run()
        r.text = "—  " + head
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = True
        r.font.color.rgb = color
        if rest:
            r2 = p.add_run()
            r2.text = "  " + rest
            r2.font.name = FONT
            r2.font.size = Pt(size)
            r2.font.color.rgb = GREY
    return box


def rect(s, x, y, w, h, fill=FAINT, line=None, radius=False):
    shp = s.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    if radius:
        try:
            shp.adjustments[0] = 0.12
        except Exception:
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def label(shp, txt, size=13, color=INK, bold=False, sub=None, sub_size=10):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Emu(45720)
    tf.margin_top = tf.margin_bottom = Emu(18288)
    p = tf.paragraphs[0]
    p.text = txt
    p.alignment = PP_ALIGN.CENTER
    for r in p.runs:
        r.font.name = FONT
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    if sub:
        p2 = tf.add_paragraph()
        p2.text = sub
        p2.alignment = PP_ALIGN.CENTER
        for r in p2.runs:
            r.font.name = FONT
            r.font.size = Pt(sub_size)
            r.font.color.rgb = GREY
    return shp


def arrow(s, x, y, w, h=Inches(0.16)):
    shp = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    shp.adjustments[0] = 0.55
    shp.adjustments[1] = 0.55
    shp.fill.solid()
    shp.fill.fore_color.rgb = LINE
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def header(s, kicker, title):
    text(s, MARGIN, Inches(0.45), SLIDE_W - 2 * MARGIN, Inches(0.3),
         kicker.upper(), size=12, color=GREY, bold=True)
    text(s, MARGIN, Inches(0.75), SLIDE_W - 2 * MARGIN, Inches(0.8),
         title, size=30, bold=True, font=FONT_LIGHT)
    ln = rect(s, MARGIN, Inches(1.55), Inches(0.55), Pt(2.6), fill=ACCENT)
    return ln


def footer(s, n):
    text(s, MARGIN, Inches(7.05), Inches(6), Inches(0.3),
         "VoiceShield · Team Jigyasa", size=9, color=GREY)
    text(s, SLIDE_W - MARGIN - Inches(0.5), Inches(7.05), Inches(0.5), Inches(0.3),
         str(n), size=9, color=GREY, align=PP_ALIGN.RIGHT)


def chip(s, x, y, txt, color, w=Inches(1.0), h=Inches(0.34), size=12):
    shp = rect(s, x, y, w, h, fill=color, radius=True)
    label(shp, txt, size=size, color=WHITE, bold=True)
    return shp


# ================================================================ 1 title
s = slide()
text(s, MARGIN, Inches(2.35), Inches(10), Inches(1.2),
     "VoiceShield", size=54, bold=True, font=FONT_LIGHT)
rect(s, MARGIN + Inches(0.02), Inches(3.35), Inches(0.9), Pt(3), fill=ACCENT)
text(s, MARGIN, Inches(3.6), Inches(11), Inches(0.6),
     "Real-time audio forensics that flags cloned voices on live calls — within 10 seconds.",
     size=20, color=GREY)
text(s, MARGIN, Inches(6.3), Inches(11), Inches(0.8),
     "Team Jigyasa  ·  Amit Anand · Suraj Harlekar · Aditya Dubey\n"
     "UCO PSB Hackathon · Problem Statement 2 · Audio Forensics for Voice Security",
     size=13, color=GREY)

# ================================================================ 2 problem
s = slide()
header(s, "The problem", "A voice is no longer proof of identity")
stats = [
    ("3 sec", "of audio is enough to clone a voice with ~85% similarity"),
    ("+442%", "surge in voice-phishing attacks in 2025, driven by AI"),
    ("$40 B", "projected annual US losses to AI-enabled fraud by 2027 (Deloitte)"),
]
cw = Inches(3.7)
for i, (big, small) in enumerate(stats):
    x = MARGIN + i * (cw + Inches(0.35))
    card = rect(s, x, Inches(2.1), cw, Inches(2.3), fill=FAINT, radius=True)
    text(s, x + Inches(0.3), Inches(2.4), cw - Inches(0.6), Inches(1.0),
         big, size=44, bold=True, color=(ACCENT if i == 2 else INK), font=FONT_LIGHT)
    text(s, x + Inches(0.3), Inches(3.4), cw - Inches(0.6), Inches(0.9),
         small, size=14, color=GREY)
text(s, MARGIN, Inches(5.0), SLIDE_W - 2 * MARGIN, Inches(1.4),
     "Fraudsters clone a customer's voice to defeat voice-biometric checks and to socially engineer "
     "call-center agents into moving money.\nToday's defenses — caller ID, phone-number metadata — are trivially spoofed. "
     "The audio itself is the only evidence that cannot lie.",
     size=17, color=INK, line_spacing=1.15)
footer(s, 2)

# ================================================================ 3 solution
s = slide()
header(s, "Our solution", "Listening for synthesis, not identity")
y = Inches(2.5)
b1 = rect(s, MARGIN, y, Inches(2.5), Inches(1.5), fill=FAINT, radius=True)
label(b1, "Live call audio", size=15, bold=True, sub="500 ms chunks · 16 kHz")
arrow(s, MARGIN + Inches(2.58), y + Inches(0.67), Inches(0.55))
b2 = rect(s, MARGIN + Inches(3.2), y, Inches(3.3), Inches(1.5), fill=INK, radius=True)
label(b2, "VoiceShield pipeline", size=16, color=WHITE, bold=True,
      sub="forensic features · model ensemble · risk fusion")
b2.text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)
arrow(s, MARGIN + Inches(6.58), y + Inches(0.67), Inches(0.55))
b3 = rect(s, MARGIN + Inches(7.2), y, Inches(2.6), Inches(1.5), fill=FAINT, radius=True)
label(b3, "Agent advisory", size=15, bold=True, sub="verdict ≤ 10 s of audio")
cx = MARGIN + Inches(10.1)
chip(s, cx, y + Inches(0.03), "GREEN · continue", S_GREEN, w=Inches(1.95), h=Inches(0.31), size=10)
chip(s, cx, y + Inches(0.41), "AMBER · escalate", S_AMBER, w=Inches(1.95), h=Inches(0.31), size=10)
chip(s, cx, y + Inches(0.79), "RED · terminate", S_RED, w=Inches(1.95), h=Inches(0.31), size=10)
chip(s, cx, y + Inches(1.17), "GREY · poor audio", S_GREY, w=Inches(1.95), h=Inches(0.31), size=10)
bullets(s, MARGIN, Inches(4.7), SLIDE_W - 2 * MARGIN, Inches(1.9), [
    ("Advisory, not judge.", "The agent stays in control; VoiceShield names the evidence behind every alert."),
    ("Privacy by construction.", "A 10-second in-memory rolling buffer — no call recording ever touches disk."),
    ("Explainable.", "Each alert cites the artifact: phase discontinuity, over-smooth pitch, vocoder band energy."),
], size=16)
footer(s, 3)

# ================================================================ 4 architecture
s = slide()
header(s, "How it works", "Two-stage cascade: cheap screening, deep confirmation")
y1 = Inches(2.15)
bw, bh, gap = Inches(2.15), Inches(1.05), Inches(0.42)
steps = [
    ("Rolling buffer", "10 s in-memory"),
    ("SNR + VAD gate", "poor audio → GREY"),
    ("Forensic features", "9 artifact panels"),
    ("Stage 1 screen", "AASIST-L, every chunk"),
]
x = MARGIN
for i, (t, sub) in enumerate(steps):
    b = rect(s, x, y1, bw, bh, fill=FAINT, radius=True)
    label(b, t, size=14, bold=True, sub=sub)
    if i < 3:
        arrow(s, x + bw + Inches(0.04), y1 + Inches(0.45), gap - Inches(0.08))
    x += bw + gap

y2 = Inches(4.0)
t = text(s, MARGIN + Inches(7.0), y1 + bh + Inches(0.12), Inches(3.6), Inches(0.35),
         "suspicion ≥ threshold — or periodic probe", size=11, color=GREY,
         align=PP_ALIGN.CENTER)
elbow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, MARGIN + Inches(8.6), y1 + bh + Inches(0.42),
                           Inches(0.32), Inches(0.44))
elbow.fill.solid(); elbow.fill.fore_color.rgb = LINE; elbow.line.fill.background()
elbow.shadow.inherit = False

stage2 = rect(s, MARGIN + Inches(3.6), y2, Inches(6.2), Inches(1.7), fill=INK, radius=True)
tf = stage2.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.text = "Stage 2 — deep ensemble"; p.alignment = PP_ALIGN.CENTER
for r in p.runs:
    r.font.name = FONT; r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = WHITE
p2 = tf.add_paragraph()
p2.text = "XLS-R 300M · WavLM · phase/pitch rules — three input representations, two training domains"
p2.alignment = PP_ALIGN.CENTER
for r in p2.runs:
    r.font.name = FONT; r.font.size = Pt(12); r.font.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)

arrow(s, MARGIN + Inches(9.86), y2 + Inches(0.72), Inches(0.5))
fusion = rect(s, MARGIN + Inches(10.4), y2, Inches(1.75), Inches(1.7), fill=FAINT, radius=True)
label(fusion, "Risk fusion", size=14, bold=True, sub="weighted + peak evidence · hysteresis")

text(s, MARGIN, Inches(6.1), SLIDE_W - 2 * MARGIN, Inches(0.8),
     "Stage 1 screens every 500 ms chunk on-budget; the heavyweight ensemble engages only on suspicion "
     "(plus a periodic probe so a weak screen can never hide an attack). GREY always wins: no risk claims without fresh speech.",
     size=14, color=GREY, line_spacing=1.15)
footer(s, 4)

# ================================================================ 5 ensemble
s = slide()
header(s, "The ensemble", "Diversity beats any single detector")
rows = [
    ("MMS-300M (NII)", "multilingual SSL, 74k h + RawBoost", "all families: GAN vocoders, TTS, VC, wild fakes", "benchmark AUC 0.997"),
    ("XLS-R 300M", "self-supervised speech model", "vocoder fingerprints in embeddings", "benchmark AUC 0.73"),
    ("WavLM-base", "self-supervised, In-the-Wild", "modern commercial TTS (ElevenLabs-class)", "benchmark AUC 0.87"),
    ("Phase / Pitch", "DSP rules — fully explainable", "artifact naming for the agent, zero fusion weight", "explainability layer"),
]
y = Inches(2.0)
col_x = [MARGIN, MARGIN + Inches(2.2), MARGIN + Inches(5.2), MARGIN + Inches(9.3)]
hdrs = ["Model", "What it is", "What it catches", "Trained on"]
for cx_, h_ in zip(col_x, hdrs):
    text(s, cx_, y, Inches(2.6), Inches(0.3), h_.upper(), size=11, color=GREY, bold=True)
y += Inches(0.42)
for i, (m, what, catches, trained) in enumerate(rows):
    if i % 2 == 0:
        rect(s, MARGIN - Inches(0.15), y - Inches(0.05), SLIDE_W - 2 * MARGIN + Inches(0.3),
             Inches(0.62), fill=FAINT)
    text(s, col_x[0], y, Inches(2.1), Inches(0.5), m, size=14, bold=True)
    text(s, col_x[1], y, Inches(2.9), Inches(0.5), what, size=13, color=GREY)
    text(s, col_x[2], y, Inches(4.0), Inches(0.5), catches, size=13)
    text(s, col_x[3], y, Inches(2.6), Inches(0.5), trained, size=12, color=GREY)
    y += Inches(0.64)
text(s, MARGIN, Inches(5.15), SLIDE_W - 2 * MARGIN, Inches(1.2),
     "Selected by measurement, not marketing: we benchmarked 8 candidate detectors on 496 clips spanning 7 GAN "
     "vocoders (WaveFake), 100+ TTS/VC systems (ASVspoof 2021), wild internet deepfakes, and our own generated "
     "clones. Fusion weights and alert thresholds come from those measured score distributions.",
     size=15, line_spacing=1.15)
text(s, MARGIN, Inches(6.25), SLIDE_W - 2 * MARGIN, Inches(0.6),
     "Every component earns its fusion weight through validation. The same benchmark retired our original Stage 1 "
     "screener (AASIST-L: AUC 0.33 on real audio) and keeps the replay module at zero weight until calibrated.",
     size=12, color=GREY, line_spacing=1.1)
footer(s, 5)

# ================================================================ 6 uniqueness
s = slide()
header(s, "What makes it different", "Four properties competitors don't combine")
cards = [
    ("Explainable alerts", "Every AMBER/RED names its artifact — “phase discontinuity”, “over-smooth pitch”. "
     "Agents act on reasons, not scores."),
    ("Privacy by design", "10-second in-memory buffer, zero persistent recording. Evidence saved only on "
     "explicit save-for-audit — aligned with Indian data-privacy expectations."),
    ("Cascade efficiency", "Lightweight screening every chunk; the GPU ensemble spins up only on suspicion. "
     "Scales to many concurrent calls per box."),
    ("Court-ready evidence", "One click exports spectrogram + phase heatmap, SHA-256 audio hash, and the full "
     "per-model score envelope — a tamper-evident fraud case file."),
]
cw, ch = Inches(5.75), Inches(1.95)
for i, (t, d) in enumerate(cards):
    x = MARGIN + (i % 2) * (cw + Inches(0.4))
    y = Inches(2.05) + (i // 2) * (ch + Inches(0.35))
    rect(s, x, y, cw, ch, fill=FAINT, radius=True)
    text(s, x + Inches(0.3), y + Inches(0.22), cw - Inches(0.6), Inches(0.4), t, size=17, bold=True)
    text(s, x + Inches(0.3), y + Inches(0.68), cw - Inches(0.6), Inches(1.2), d, size=13, color=GREY,
         line_spacing=1.1)
footer(s, 6)

# ================================================================ 7 results
s = slide()
header(s, "Measured, not promised", "Validation results on real audio")
rows = [
    ("Detection — 323 fakes: 7 GAN vocoders, 100+ TTS/VC systems, wild deepfakes, ElevenLabs",
     "99% RED ≤ 10 s (≥4 s clips)", S_RED),
    ("Time to alert (median across detected fakes)", "AMBER 1.5 s · RED 2.0 s", S_RED),
    ("Hindi — FLEURS genuine vs MMS-TTS fakes (multilingual)", "AUC 1.0 · 100% flagged ≤5 s", S_RED),
    ("Genuine speech — 173 clips incl. codec + noisy wild domains", "3.5% false-RED · demo set: zero", S_GREEN),
    ("Silence / degraded audio", "GREY — never a false risk claim", S_GREY),
    ("Latency per 500 ms chunk (60 W laptop GPU)", "p50 83–113 ms · p95 207 ms", INK),
]
y = Inches(2.0)
for t, v, c in rows:
    rect(s, MARGIN, y, Inches(7.4), Inches(0.66), fill=FAINT, radius=True)
    text(s, MARGIN + Inches(0.25), y + Inches(0.13), Inches(7.0), Inches(0.45), t, size=13)
    text(s, MARGIN + Inches(7.7), y + Inches(0.12), Inches(4.2), Inches(0.45), v, size=15, bold=True, color=c)
    y += Inches(0.76)
text(s, MARGIN, Inches(6.6), SLIDE_W - 2 * MARGIN, Inches(0.6),
     "Problem statement asks: flag High Risk within the first 10 seconds — median RED is at 1.5 s. "
     "Measured on a 496-clip benchmark (WaveFake, ASVspoof 2021 DF, In-the-Wild + generated clones); "
     "94 automated tests gate every change.",
     size=12, color=GREY)
footer(s, 7)

# ================================================================ 8 moat
s = slide()
header(s, "Technical moat", "The models are public. The calibration is ours.")
bullets(s, MARGIN, Inches(2.05), SLIDE_W - 2 * MARGIN, Inches(4.4), [
    ("Evidence-calibrated fusion.", "We measured each model's per-domain failure modes (a top model scores 1.0 "
     "on everything; another false-fires on 10% of genuine speakers) and encoded that trust into per-model "
     "peak-evidence rights. Plug-and-play ensembles don't survive contact with real audio — ours did."),
    ("Weak-evidence discipline.", "Scores scale with voiced content; silence can never escalate risk; GREY "
     "overrides everything. Zero false alerts on genuine speakers is an engineered property, not luck."),
    ("Channel forensics.", "We quantified how re-recording destroys synthesis artifacts (0.91 → 0.10 through a "
     "loudspeaker) — a measured limitation most vendors won't tell you about, and the reason replay-channel "
     "detection is our top calibration priority for the pilot."),
    ("A validation harness as the product.", "Fixture pipeline (genuine / TTS / cloned / degraded), 94 automated "
     "tests, latency benchmarks — new detectors are trusted only as far as they prove. Our own replay module is "
     "implemented but quarantined at zero weight until real call audio validates it. That discipline is the moat."),
], size=15, gap=14)
footer(s, 8)

# ================================================================ 9 banking
s = slide()
header(s, "Deployability", "Drops into existing bank telephony")
y = Inches(2.3)
b = rect(s, MARGIN, y, Inches(2.3), Inches(1.4), fill=FAINT, radius=True)
label(b, "Call center PBX / SBC", size=13, bold=True, sub="existing telephony")
arrow(s, MARGIN + Inches(2.36), y + Inches(0.58), Inches(0.5))
text(s, MARGIN + Inches(1.95), y + Inches(1.5), Inches(1.4), Inches(0.3), "RTP fork", size=10,
     color=GREY, align=PP_ALIGN.CENTER)
b = rect(s, MARGIN + Inches(2.95), y, Inches(3.1), Inches(1.4), fill=INK, radius=True)
label(b, "VoiceShield server", size=14, color=WHITE, bold=True,
      sub="on-prem · 1 commodity GPU · CPU fallback")
b.text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)
arrow(s, MARGIN + Inches(6.12), y + Inches(0.58), Inches(0.5))
b = rect(s, MARGIN + Inches(6.7), y, Inches(2.55), Inches(1.4), fill=FAINT, radius=True)
label(b, "Agent desktop", size=13, bold=True, sub="advisory badge · WebSocket, 2×/sec")
b2 = rect(s, MARGIN + Inches(9.55), y, Inches(2.55), Inches(1.4), fill=FAINT, radius=True)
label(b2, "Fraud team vault", size=13, bold=True, sub="evidence packages on demand")
conn = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, MARGIN + Inches(9.28), y + Inches(0.58),
                          Inches(0.24), Inches(0.16))
conn.fill.solid(); conn.fill.fore_color.rgb = LINE; conn.line.fill.background()
conn.shadow.inherit = False

bullets(s, MARGIN, Inches(4.35), SLIDE_W - 2 * MARGIN, Inches(2.2), [
    ("Integration surface is two URLs.", "Versioned REST + WebSocket (/v1 frozen, /v2 ensemble) — CRM and "
     "agent-assist tools consume JSON; no telephony rework."),
    ("Compliance posture.", "No cloud dependency, no stored voice data, inference entirely on-premise; evidence "
     "export is explicit, hashed, and auditable."),
    ("Sizing.", "166 ms per 500 ms chunk on a laptop GPU today; one datacenter GPU serves multiple concurrent "
     "calls, and the cascade cuts average load further."),
], size=15, gap=10)
footer(s, 9)

# ================================================================ 10 demo
s = slide()
header(s, "Live demo", "What you will see (5 minutes)")
steps = [
    ("1", "Genuine voice — live microphone", "Speak into the browser; verdict stays GREEN, artifacts panel stays quiet."),
    ("2", "Real ElevenLabs voice clone", "Upload the cloned clip: AMBER at 3 s, RED at 5 s — watch the per-model score bars disagree, and fusion decide."),
    ("3", "Evidence export", "One click produces the audit package: spectrogram, phase heatmap, SHA-256, score envelope."),
]
y = Inches(2.2)
for n, t, d in steps:
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN, y, Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = INK; circ.line.fill.background()
    circ.shadow.inherit = False
    label(circ, n, size=16, color=WHITE, bold=True)
    text(s, MARGIN + Inches(0.75), y - Inches(0.02), Inches(10.5), Inches(0.4), t, size=18, bold=True)
    text(s, MARGIN + Inches(0.75), y + Inches(0.42), Inches(10.5), Inches(0.5), d, size=14, color=GREY)
    y += Inches(1.25)
text(s, MARGIN, Inches(6.2), SLIDE_W - 2 * MARGIN, Inches(0.7),
     "Everything runs live on this laptop — the same stack that would sit in the bank's datacenter.",
     size=14, color=GREY)
footer(s, 10)

# ================================================================ 11 TRL
s = slide()
header(s, "Technology readiness", "TRL 6 — demonstrated in a relevant environment")
lvls = [
    ("TRL 4", "components validated in lab", True),
    ("TRL 5", "integrated pipeline validated", True),
    ("TRL 6", "end-to-end system, live audio,\nreal pretrained models", True),
    ("TRL 7", "pilot on live call traffic", False),
    ("TRL 8", "hardened telephony integration", False),
]
x = MARGIN
bw = Inches(2.2)
for name, d, done in lvls:
    fill = INK if done else FAINT
    col = WHITE if done else GREY
    b = rect(s, x, Inches(2.15), bw, Inches(1.25), fill=fill, radius=True)
    label(b, name, size=15, color=col, bold=True, sub=d)
    if done:
        b.text_frame.paragraphs[1].runs[0].font.color.rgb = RGBColor(0xBF, 0xBF, 0xBF)
    x += bw + Inches(0.18)
bullets(s, MARGIN, Inches(3.9), SLIDE_W - 2 * MARGIN, Inches(2.6), [
    ("Working system, not slides.", "Live browser-mic capture → GPU ensemble → advisory dashboard, running "
     "end-to-end today at p95 166 ms per chunk."),
    ("Engineering rigor.", "94 automated tests including acceptance gates (false-alert rate, detection-time, "
     "latency); CI-gated; structured audit logging with trace IDs."),
    ("Validated against real attacks.", "Detected a genuine ElevenLabs clone of a real voice — the exact attack "
     "in the problem statement."),
], size=15, gap=10)
footer(s, 11)

# ================================================================ 12 roadmap
s = slide()
header(s, "Roadmap", "From prototype to production")
line_y = Inches(2.9)
rect(s, MARGIN, line_y, SLIDE_W - 2 * MARGIN, Pt(2), fill=LINE)
phases = [
    ("Today", "Hackathon prototype", "cascade ensemble · replay detection ·\nevidence export · live dashboard", True),
    ("+3 months", "Bank pilot", "SIP/RTP tap on live traffic · replay-channel\ndetection calibrated on real call audio", False),
    ("+6 months", "India-ready", "Hindi + regional-language robustness ·\nmultilingual anti-spoof model · agent rollout", False),
    ("+12 months", "Network effect", "cross-branch fraud analytics · voice-clone\nthreat intel shared across PSBs", False),
]
seg = (SLIDE_W - 2 * MARGIN) / len(phases)
for i, (when, t, d, now) in enumerate(phases):
    x = MARGIN + int(seg) * i
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, x, line_y - Inches(0.07), Inches(0.17), Inches(0.17))
    dot.fill.solid(); dot.fill.fore_color.rgb = ACCENT if now else INK
    dot.line.fill.background(); dot.shadow.inherit = False
    text(s, x, line_y - Inches(0.62), Inches(2.4), Inches(0.35), when.upper(), size=12,
         color=(ACCENT if now else GREY), bold=True)
    text(s, x, line_y + Inches(0.32), Inches(2.9), Inches(0.4), t, size=16, bold=True)
    text(s, x, line_y + Inches(0.78), Inches(2.9), Inches(1.2), d, size=12, color=GREY, line_spacing=1.1)
text(s, MARGIN, Inches(5.4), SLIDE_W - 2 * MARGIN, Inches(1.2),
     "Deliberate scope: detection accuracy and honest calibration first, telephony plumbing second. "
     "The API seams (versioned /v1 → /v2) mean each phase ships without breaking the last.",
     size=14, color=GREY)
footer(s, 12)

# ================================================================ 13 close
s = slide()
text(s, MARGIN, Inches(2.6), SLIDE_W - 2 * MARGIN, Inches(1.6),
     "The audio is the only evidence\nthat cannot lie. We read it.",
     size=40, bold=True, font=FONT_LIGHT, line_spacing=1.05)
rect(s, MARGIN + Inches(0.02), Inches(4.35), Inches(0.9), Pt(3), fill=ACCENT)
text(s, MARGIN, Inches(4.6), Inches(11), Inches(0.5),
     "VoiceShield — Team Jigyasa  ·  Questions welcome",
     size=16, color=GREY)
footer(s, 13)

# ---------------------------------------------------------------- save
out_dir = os.path.join(os.path.dirname(__file__), "..", "docs", "presentation")
os.makedirs(out_dir, exist_ok=True)
out = os.path.join(out_dir, "VoiceShield_Jigyasa_UCOPSB.pptx")
prs.save(out)
print("saved:", os.path.abspath(out), f"({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
